# -*- coding: utf-8 -*-
"""덱 레이아웃 근사 렌더러 (QA용). pptx 도형을 읽어 contact sheet PNG 생성.
정확한 렌더가 아니라 위치/오버플로/이미지배치 점검용."""
import io, math, os
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
from matplotlib.patches import FancyBboxPatch, Rectangle
from matplotlib import font_manager
from PIL import Image
from pptx import Presentation
from pptx.util import Emu
from pptx.enum.shapes import MSO_SHAPE_TYPE

for cand in ["AppleGothic", "Apple SD Gothic Neo"]:
    try:
        font_manager.findfont(cand, fallback_to_default=False)
        plt.rcParams["font.family"] = cand; break
    except Exception:
        pass
plt.rcParams["axes.unicode_minus"] = False

HERE = os.path.dirname(os.path.abspath(__file__))
prs = Presentation(os.path.join(HERE, "아보하_HCI_최종발표.pptx"))
SW = prs.slide_width / 914400.0
SH = prs.slide_height / 914400.0

def emu_in(v):
    return (v or 0) / 914400.0

def rgb_of(shape):
    try:
        if shape.fill.type is not None and int(shape.fill.type) == 1:  # solid
            c = shape.fill.fore_color.rgb
            return "#%02x%02x%02x" % (c[0], c[1], c[2])
    except Exception:
        pass
    try:  # gradient → first stop
        gs = shape.fill.gradient_stops
        c = gs[0].color.rgb
        return "#%02x%02x%02x" % (c[0], c[1], c[2])
    except Exception:
        return None

def run_color(run):
    try:
        c = run.font.color.rgb
        return "#%02x%02x%02x" % (c[0], c[1], c[2])
    except Exception:
        return "#1b1c26"

overflows = []
SC = 1.0  # inches
def render_slide(slide, idx, ax):
    ax.set_xlim(0, SW); ax.set_ylim(0, SH); ax.invert_yaxis()
    ax.set_xticks([]); ax.set_yticks([])
    for sp in ("top","bottom","left","right"): ax.spines[sp].set_color("#cccccc")
    for sh in slide.shapes:
        L, T, W, H = emu_in(sh.left), emu_in(sh.top), emu_in(sh.width), emu_in(sh.height)
        # overflow check
        if L < -0.05 or T < -0.05 or L+W > SW+0.05 or T+H > SH+0.05:
            overflows.append((idx, sh.shape_type, round(L,2), round(T,2), round(W,2), round(H,2)))
        if sh.shape_type == MSO_SHAPE_TYPE.PICTURE:
            try:
                im = Image.open(io.BytesIO(sh.image.blob)).convert("RGBA")
                ax.imshow(im, extent=[L, L+W, T+H, T], zorder=3, aspect="auto", interpolation="bilinear")
            except Exception:
                ax.add_patch(Rectangle((L,T), W, H, facecolor="#dddddd", edgecolor="none", zorder=3))
            continue
        fill = rgb_of(sh)
        if fill and W > 0 and H > 0:
            ax.add_patch(Rectangle((L,T), W, H, facecolor=fill, edgecolor="none", zorder=2))
        if sh.has_text_frame and sh.text_frame.text.strip():
            tf = sh.text_frame
            yoff = T + 0.18
            for p in tf.paragraphs:
                if not p.runs:
                    continue
                txt = "".join(r.text for r in p.runs)
                if not txt.strip():
                    yoff += 0.18; continue
                r0 = p.runs[0]
                sz = (r0.font.size.pt if r0.font.size else 14)
                col = run_color(r0)
                bold = bool(r0.font.bold)
                # rough wrap by width
                maxchars = max(6, int(W / (sz*0.0125)))
                line = txt if len(txt) <= maxchars else txt[:maxchars-1] + "…"
                ax.text(L+0.12, yoff, line, fontsize=max(5, sz*0.62), color=col,
                        fontweight=("bold" if bold else "normal"), va="top", zorder=5)
                yoff += max(0.2, sz*0.013)
                if yoff > T + H + 0.05:
                    break
    ax.text(0.05, 0.05, f"S{idx}", fontsize=7, color="#999", va="top", transform=ax.transAxes)

n = len(prs.slides)
cols, rows = 3, math.ceil(n/3)
fig, axes = plt.subplots(rows, cols, figsize=(cols*4.2, rows*2.45))
axes = axes.flatten()
for i, slide in enumerate(prs.slides):
    render_slide(slide, i, axes[i])
for j in range(n, len(axes)):
    axes[j].axis("off")
plt.tight_layout(pad=0.6)
out = os.path.join(HERE, "preview_contact.png")
fig.savefig(out, dpi=110, bbox_inches="tight")
print("preview ->", out)
print("OVERFLOWS:", len(overflows))
for o in overflows:
    print("  slide", o[0], "type", o[1], "L,T,W,H=", o[2:])
