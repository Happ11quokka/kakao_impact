# -*- coding: utf-8 -*-
"""
아보하(하루보석) — HCI Project II 최종 발표 덱 빌더
- matplotlib 차트 생성 → charts/
- python-pptx 로 16:9 덱 생성 → 아보하_HCI_최종발표.pptx
- 슬라이드별 발표자 노트 + outline.md 자동 생성
실행: python3 build_deck.py
"""
import os
from PIL import Image

HERE = os.path.dirname(os.path.abspath(__file__))
CHARTS = os.path.join(HERE, "charts")
REPO = os.path.abspath(os.path.join(HERE, "..", ".."))
ASSET = os.path.join(REPO, "2_avoha")
os.makedirs(CHARTS, exist_ok=True)

def asset(*p):
    return os.path.join(ASSET, *p)

# ---------- 브랜드 팔레트 (design/brand/tokens.css) ----------
HEX = dict(
    coral="#E8614D", mint="#3AAFA9", amber="#E8A838",
    dawn="#F4EBD9", dusk="#3A3E5B", ink="#1B1C26", parch="#FBF7EE",
    cream="#FFFAF4", sunshine="#F4A836", peach="#FDE8C8", brown="#5A3E28",
    beige="#F0E0C8", slate="#5B6079", white="#FFFFFF", muted="#8A8FA3",
    good="#36B37E", bad="#E5573F",
)

# =====================================================================
# 1) 차트 생성 (matplotlib)
# =====================================================================
def make_charts():
    import matplotlib
    matplotlib.use("Agg")
    import matplotlib.pyplot as plt
    from matplotlib import font_manager
    # 한글 폰트
    for cand in ["AppleGothic", "Apple SD Gothic Neo", "NanumGothic", "Noto Sans CJK KR"]:
        try:
            font_manager.findfont(cand, fallback_to_default=False)
            plt.rcParams["font.family"] = cand
            break
        except Exception:
            continue
    plt.rcParams["axes.unicode_minus"] = False
    C = HEX
    def save(fig, name):
        fig.savefig(os.path.join(CHARTS, name), dpi=200, bbox_inches="tight",
                    transparent=True)
        plt.close(fig)
    def style(ax):
        for s in ("top", "right"):
            ax.spines[s].set_visible(False)
        ax.spines["left"].set_color("#CBC8BE"); ax.spines["bottom"].set_color("#CBC8BE")
        ax.tick_params(colors=C["ink"], labelsize=12)
        ax.set_axisbelow(True)

    # 1. 감정인지 사전/사후
    fig, ax = plt.subplots(figsize=(6.6, 4.3))
    bars = ax.bar(["사전(n=8)", "사후(n=8)"], [3.60, 3.90],
                  color=[C["slate"], C["coral"]], width=0.55, zorder=3)
    ax.set_ylim(0, 5); ax.set_ylabel("감정인지 자기보고 (5점)", fontsize=12)
    for b, v in zip(bars, [3.60, 3.90]):
        ax.text(b.get_x()+b.get_width()/2, v+0.08, f"{v:.2f}", ha="center",
                fontsize=15, fontweight="bold", color=C["ink"])
    ax.annotate("", xy=(1, 4.25), xytext=(0, 3.95),
                arrowprops=dict(arrowstyle="->", color=C["coral"], lw=2))
    ax.text(0.5, 4.55, "+0.30", ha="center", color=C["coral"], fontsize=14, fontweight="bold")
    ax.grid(axis="y", color="#E8E4DA", zorder=0)
    style(ax); save(fig, "emotion_recognition.png")

    # 2. 정확도 Before/After (사용자 동의율)
    fig, ax = plt.subplots(figsize=(6.6, 4.3))
    bars = ax.bar(["개선 직전\n(05.13–20)", "2차 MVP\n(05.21–28)"], [41.7, 97.3],
                  color=[C["muted"], C["good"]], width=0.55, zorder=3)
    ax.set_ylim(0, 110); ax.set_ylabel("AI 분류 사용자 동의율 (%)", fontsize=12)
    for b, v in zip(bars, [41.7, 97.3]):
        ax.text(b.get_x()+b.get_width()/2, v+2, f"{v:.1f}%", ha="center",
                fontsize=15, fontweight="bold", color=C["ink"])
    ax.grid(axis="y", color="#E8E4DA", zorder=0)
    style(ax); save(fig, "accuracy_beforeafter.png")

    # 3. SUS
    fig, ax = plt.subplots(figsize=(6.6, 4.3))
    b = ax.bar(["아보하 SUS"], [71.25], color=C["mint"], width=0.4, zorder=3,
               yerr=[[11.24], [11.24]], capsize=8, ecolor=C["slate"])
    ax.axhline(68, color=C["coral"], ls="--", lw=2, zorder=2)
    ax.text(0.42, 68, "업계 평균 68", color=C["coral"], fontsize=12, va="bottom")
    ax.text(0, 71.25+13, "71.25", ha="center", fontsize=16, fontweight="bold", color=C["ink"])
    ax.set_ylim(0, 100); ax.set_ylabel("SUS 점수 (0–100)", fontsize=12)
    ax.text(0, 5, "95% CI 60–82", ha="center", fontsize=11, color=C["muted"])
    ax.grid(axis="y", color="#E8E4DA", zorder=0)
    style(ax); save(fig, "sus.png")

    # 4. 전환율 간극 (자기보고 vs 행동)
    fig, ax = plt.subplots(figsize=(6.6, 4.3))
    bars = ax.bar(["설문 자기보고", "행동 로그"], [93, 24],
                  color=[C["amber"], C["slate"]], width=0.55, zorder=3)
    ax.set_ylim(0, 110); ax.set_ylabel("카카오톡 → 웹 전환율 (%)", fontsize=12)
    for b, v, t in zip(bars, [93, 24], ["93%", "≈24%"]):
        ax.text(b.get_x()+b.get_width()/2, v+2, t, ha="center",
                fontsize=15, fontweight="bold", color=C["ink"])
    ax.annotate("", xy=(1, 60), xytext=(0, 60),
                arrowprops=dict(arrowstyle="<->", color=C["coral"], lw=1.8))
    ax.text(0.5, 63, "자기보고–행동 간극", ha="center", color=C["coral"], fontsize=12, fontweight="bold")
    ax.grid(axis="y", color="#E8E4DA", zorder=0)
    style(ax); save(fig, "conversion_gap.png")

    # 5. 회고 기능 만족도
    fig, ax = plt.subplots(figsize=(6.8, 4.3))
    labels = ["캘린더", "감정분석 질문", "Recap", "원석·로기", "원석 흐름 이해"]
    vals = [3.77, 3.69, 3.54, 3.46, 3.38]
    cols = [C["coral"]] + [C["mint"]]*3 + [C["muted"]]
    y = range(len(labels))[::-1]
    bars = ax.barh(list(y), vals, color=cols, height=0.62, zorder=3)
    ax.set_yticks(list(y)); ax.set_yticklabels(labels, fontsize=12)
    ax.set_xlim(0, 5); ax.set_xlabel("5점 만점", fontsize=11)
    for b, v in zip(bars, vals):
        ax.text(v+0.06, b.get_y()+b.get_height()/2, f"{v:.2f}", va="center",
                fontsize=12, fontweight="bold", color=C["ink"])
    ax.grid(axis="x", color="#E8E4DA", zorder=0)
    style(ax); save(fig, "recap_features.png")

    # 6. 이탈 원인
    fig, ax = plt.subplots(figsize=(6.6, 4.3))
    labels = ["잊어버림", "귀찮음", "사용 이유 부족"]
    vals = [11, 10, 6]
    bars = ax.bar(labels, vals, color=[C["coral"], C["amber"], C["slate"]], width=0.6, zorder=3)
    for b, v in zip(bars, vals):
        ax.text(b.get_x()+b.get_width()/2, v+0.15, str(v), ha="center",
                fontsize=14, fontweight="bold", color=C["ink"])
    ax.set_ylabel("응답 수 (복수 선택)", fontsize=12); ax.set_ylim(0, 13)
    ax.grid(axis="y", color="#E8E4DA", zorder=0)
    style(ax); save(fig, "churn_reasons.png")

    # 7. NPS 도넛
    fig, ax = plt.subplots(figsize=(5.4, 4.3))
    sizes = [21.4, 28.6, 50.0]
    cols = [C["good"], C["amber"], C["bad"]]
    wedges, _ = ax.pie(sizes, colors=cols, startangle=90, counterclock=False,
                       wedgeprops=dict(width=0.42, edgecolor="white", linewidth=2))
    ax.text(0, 0.12, "NPS", ha="center", fontsize=13, color=C["muted"])
    ax.text(0, -0.16, "-28.6", ha="center", fontsize=24, fontweight="bold", color=C["bad"])
    ax.legend(wedges, ["추천 21.4%", "중립 28.6%", "비추천 50.0%"],
              loc="lower center", bbox_to_anchor=(0.5, -0.12), ncol=3, fontsize=9,
              frameon=False)
    save(fig, "nps.png")

    # 8. 리텐션 (활동일수 분포, 활동자 29명)
    fig, ax = plt.subplots(figsize=(6.6, 4.3))
    labels = ["1일", "2일", "3일", "4일", "5–8일", "9일+"]
    vals = [41.4, 20.7, 13.8, 6.9, 0.0, 17.2]
    cols = [C["coral"]] + [C["slate"]]*4 + [C["mint"]]
    bars = ax.bar(labels, vals, color=cols, width=0.62, zorder=3)
    for b, v in zip(bars, vals):
        if v > 0:
            ax.text(b.get_x()+b.get_width()/2, v+0.6, f"{v:.0f}%", ha="center",
                    fontsize=12, fontweight="bold", color=C["ink"])
    ax.set_ylabel("활동 사용자 비율 (n=29)", fontsize=12); ax.set_ylim(0, 48)
    ax.grid(axis="y", color="#E8E4DA", zorder=0)
    style(ax); save(fig, "retention.png")

    # 9. 클래스별 정밀도 (annoyance 약점)
    fig, ax = plt.subplots(figsize=(6.6, 4.3))
    labels = ["평온", "무탈", "슬픔", "기쁨", "위로", "짜증"]
    vals = [100, 100, 90, 78, 78, 50]
    cols = [C["mint"]]*5 + [C["coral"]]
    bars = ax.bar(labels, vals, color=cols, width=0.62, zorder=3)
    for b, v in zip(bars, vals):
        ax.text(b.get_x()+b.get_width()/2, v+1.5, f"{v:.0f}", ha="center",
                fontsize=12, fontweight="bold", color=C["ink"])
    ax.set_ylabel("AI 예측 정밀도 (%)", fontsize=12); ax.set_ylim(0, 112)
    ax.grid(axis="y", color="#E8E4DA", zorder=0)
    style(ax); save(fig, "precision_class.png")

    print("charts: done ->", CHARTS)


# 차트를 먼저 생성한다 (덱이 PNG를 참조하므로 반드시 선행)
make_charts()

# =====================================================================
# 2) 덱 빌더 (python-pptx)
# =====================================================================
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

def rgb(h):
    h = h.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))

C = {k: rgb(v) for k, v in HEX.items()}
KFONT = "Apple SD Gothic Neo"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]
OUTLINE = []  # for outline.md


def _set_font(run, name, size, color, bold):
    run.font.size = Pt(size); run.font.bold = bold
    run.font.color.rgb = color; run.font.name = name
    rPr = run._r.get_or_add_rPr()
    for tag in ("a:latin", "a:ea", "a:cs"):
        el = rPr.find(qn(tag))
        if el is None:
            el = rPr.makeelement(qn(tag), {}); rPr.append(el)
        el.set("typeface", name)


def R(t, size, color, bold=False, font=None):
    return dict(t=t, size=size, color=color, bold=bold, font=font or KFONT)

def P(runs, sa=6, sb=0, line=None, align=None):
    return dict(runs=runs, sa=sa, sb=sb, line=line, align=align)


def slide():
    s = prs.slides.add_slide(BLANK)
    return s

def bg(s, *colors, angle=90):
    sp = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    sp.line.fill.background(); sp.shadow.inherit = False
    if len(colors) == 1:
        sp.fill.solid(); sp.fill.fore_color.rgb = colors[0]
    else:
        sp.fill.gradient()
        stops = sp.fill.gradient_stops
        stops[0].position = 0.0; stops[0].color.rgb = colors[0]
        stops[-1].position = 1.0; stops[-1].color.rgb = colors[-1]
        try: sp.fill.gradient_angle = angle
        except Exception: pass
    return sp

def rect(s, x, y, w, h, fill=None, line=None, lw=1.0, rounded=False, radius=0.08, shadow=False):
    shape = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    sp = s.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    sp.shadow.inherit = False
    if rounded:
        try: sp.adjustments[0] = radius
        except Exception: pass
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line; sp.line.width = Pt(lw)
    if shadow:
        el = sp._element.spPr
        from pptx.oxml.ns import qn as _q
        ef = el.makeelement(_q("a:effectLst"), {})
        sh = el.makeelement(_q("a:outerShdw"),
                            {"blurRad": "90000", "dist": "40000", "dir": "5400000", "rotWithShape": "0"})
        clr = el.makeelement(_q("a:srgbClr"), {"val": "1B1C26"})
        alpha = el.makeelement(_q("a:alpha"), {"val": "22000"})
        clr.append(alpha); sh.append(clr); ef.append(sh); el.append(ef)
    return sp

def txt(s, x, y, w, h, paras, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, wrap=True):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = wrap; tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    for i, pa in enumerate(paras):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = pa.get("align") or align
        p.space_after = Pt(pa.get("sa", 6)); p.space_before = Pt(pa.get("sb", 0))
        if pa.get("line"):
            p.line_spacing = pa["line"]
        for r in pa["runs"]:
            run = p.add_run(); run.text = r["t"]
            _set_font(run, r["font"], r["size"], r["color"], r["bold"])
    return tb

def bullets(s, x, y, w, h, items, size=15, color=None, gap=8, mark="•  ", markcolor=None, line=1.15):
    color = color or C["ink"]; markcolor = markcolor or C["coral"]
    paras = []
    for it in items:
        if isinstance(it, tuple):
            label, sub = it
            paras.append(P([R(mark, size, markcolor, True), R(label, size, color, True)], sa=2, line=line))
            paras.append(P([R("     " + sub, size-2, C["slate"])], sa=gap, line=line))
        else:
            paras.append(P([R(mark, size, markcolor, True), R(it, size, color)], sa=gap, line=line))
    return txt(s, x, y, w, h, paras)

def pic_fit(s, path, x, y, w, h, ha="center", va="middle"):
    if not os.path.exists(path):
        return None
    iw, ih = Image.open(path).size
    ar = iw / ih
    if ar > w / h:
        nw, nh = w, w / ar
    else:
        nh, nw = h, h * ar
    px = x + (w - nw) * (0.5 if ha == "center" else (1.0 if ha == "right" else 0.0))
    py = y + (h - nh) * (0.5 if va == "middle" else (1.0 if va == "bottom" else 0.0))
    return s.shapes.add_picture(path, Inches(px), Inches(py), Inches(nw), Inches(nh))

def pill(s, x, y, label, fill, tcolor, size=12, padw=0.16):
    w = 0.16 * len(label) + padw * 2 + 0.1
    sp = rect(s, x, y, w, 0.42, fill=fill, rounded=True, radius=0.5)
    txt(s, x, y, w, 0.42, [P([R(label, size, tcolor, True)], sa=0)],
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    return w

def note(s, text):
    s.notes_slide.notes_text_frame.text = text

def record(n, title, bullets_list, spk):
    OUTLINE.append((n, title, bullets_list, spk))

MARGIN = 0.7
CW = 13.333 - MARGIN * 2  # content width

def kicker_title(s, kicker, title, tnum, color=None, tcolor=None, y=0.62):
    color = color or C["coral"]; tcolor = tcolor or C["ink"]
    rect(s, MARGIN, y+0.02, 0.12, 0.62, fill=color)
    txt(s, MARGIN+0.28, y, CW-0.28, 0.4,
        [P([R(kicker, 13, color, True)], sa=0)])
    txt(s, MARGIN+0.28, y+0.34, CW-0.28, 0.7,
        [P([R(title, 27, tcolor, True)], sa=0)])
    txt(s, 13.333-MARGIN-0.6, y+0.05, 0.6, 0.4,
        [P([R(tnum, 12, C["muted"], True)], sa=0)], align=PP_ALIGN.RIGHT)

def footer(s, dark=False):
    col = C["beige"] if dark else C["muted"]
    txt(s, MARGIN, 7.04, 8, 0.3, [P([R("아보하 · 하루보석  |  HCI Project II", 9, col)], sa=0)])


# ---------------------------------------------------------------------
# SLIDE 0 — 표지
# ---------------------------------------------------------------------
s = slide()
bg(s, C["dusk"], C["ink"], angle=120)
# 별 점 장식
import random
random.seed(7)
for _ in range(44):
    rx, ry = random.uniform(0.3, 12.9), random.uniform(0.3, 4.6)
    d = random.choice([0.03, 0.045, 0.06])
    o = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(rx), Inches(ry), Inches(d), Inches(d))
    o.fill.solid(); o.fill.fore_color.rgb = C["beige"]; o.line.fill.background(); o.shadow.inherit = False
# 달
moon = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(11.7), Inches(0.7), Inches(0.7), Inches(0.7))
moon.fill.solid(); moon.fill.fore_color.rgb = C["amber"]; moon.line.fill.background(); moon.shadow.inherit = False
pic_fit(s, asset("frontend", "public", "images", "mascot.png"), 9.7, 3.7, 3.0, 3.0, ha="center", va="bottom")
txt(s, MARGIN, 2.35, 10, 0.5, [P([R("HCI Project II · 최종 발표", 15, C["amber"], True)], sa=0)])
txt(s, MARGIN, 2.95, 11, 1.6, [
    P([R("아보하", 52, C["white"], True), R("  하루보석", 30, C["beige"], True)], sa=2),
    P([R("카카오톡으로 ‘채집’하고 웹에서 돌아보는 청년 감정인지 솔루션", 18, C["beige"])], sb=8),
])
txt(s, MARGIN, 4.95, 11, 0.5,
    [P([R("기록을 채집으로, 감정을 원석으로 — 설계 · 실험 · 검증의 HCI 여정", 14, C["mint"], True)], sa=0)])
txt(s, MARGIN, 6.5, 11, 0.5,
    [P([R("닥토 공방 · 임동현 · 한양대학교 인간컴퓨터상호작용", 12, C["beige"])], sa=0)])
note(s, "안녕하세요. 저희 팀 프로젝트 ‘아보하(하루보석)’의 최종 발표를 시작하겠습니다. "
        "아보하는 ‘아주 보통의 하루’에서도 흘려보내던 감정을 알아차리게 돕는 서비스입니다. "
        "오늘은 기능 소개가 아니라, 어떤 가설을 세우고 어떻게 실험·검증했는지 HCI 관점에서 보여드리겠습니다.")
record(0, "표지 — 아보하(하루보석)", ["부제: 카카오톡 기반 청년 감정인지 솔루션", "HCI 관점의 설계·실험·검증"],
       "기능 소개가 아니라 가설–실험–검증의 HCI 여정임을 예고.")

# ---------------------------------------------------------------------
# SLIDE 1 — 목차
# ---------------------------------------------------------------------
s = slide()
bg(s, C["parch"])
kicker_title(s, "AGENDA", "오늘의 흐름", "01")
items = [
    ("Ⅰ. 문제정의 (Idea)", "왜 청년의 ‘감정인지’인가"),
    ("Ⅱ. 솔루션 & 프로토타입", "채집–원석–회고 / HCI 설계 근거"),
    ("Ⅲ. 실험 설계", "가설 · KPI · 척도 · 혼합방법"),
    ("Ⅳ. 결과 (User Study)", "설문 · 행동 · 인터뷰"),
    ("Ⅴ. 고찰", "무엇을 배웠나 · 한계 · 이론"),
    ("Ⅵ. 향후 — 통제 A/B 설계", "다음 실험과 비전"),
]
gx, gy, gw, gh = MARGIN, 2.1, (CW-0.6)/2, 1.25
for i, (t, sub) in enumerate(items):
    col = gx + (i % 2) * (gw + 0.6)
    row = gy + (i // 2) * (gh + 0.18)
    rect(s, col, row, gw, gh, fill=C["white"], rounded=True, radius=0.1, shadow=True)
    rect(s, col, row, 0.14, gh, fill=[C["coral"], C["amber"], C["mint"]][i // 2], rounded=False)
    txt(s, col+0.34, row+0.22, gw-0.5, 0.5, [P([R(t, 16, C["ink"], True)], sa=0)])
    txt(s, col+0.34, row+0.72, gw-0.5, 0.4, [P([R(sub, 12, C["slate"])], sa=0)])
footer(s)
note(s, "발표는 여섯 단계입니다. 문제정의에서 시작해 솔루션과 프로토타입, 그리고 핵심인 실험 설계와 "
        "사용자 스터디 결과, 고찰, 마지막으로 다음 통제 실험 설계까지 다루겠습니다.")
record(1, "목차", [t for t, _ in items], "6단계 흐름 안내. 실험 설계와 결과가 중심임을 강조.")

# ---------------------------------------------------------------------
# SLIDE 2 — 문제정의 (1) 청년 마음건강
# ---------------------------------------------------------------------
s = slide()
bg(s, C["parch"])
kicker_title(s, "Ⅰ. 문제정의 — IDEA", "청년 마음건강 문제가 구조적으로 심화되고 있다", "02")
stats = [
    ("19.4%", "우울증 진료 1위 = 20대", "’18→’22 +90% (심평원 2024)"),
    ("34.8%", "25–29세 번아웃 경험률", "청년 전체 32.2% (국가데이터처 2025)"),
    ("73.6%", "정신건강 문제 경험률", "전년 대비 +9.7%p (N=3,000)"),
    ("16.2%", "청년 정신건강서비스 이용률", "필요–이용의 큰 간극 (복지부 2024)"),
]
cw2 = (CW - 0.6) / 4
for i, (big, lab, sub) in enumerate(stats):
    x = MARGIN + i * (cw2 + 0.2)
    rect(s, x, 2.25, cw2, 2.5, fill=C["white"], rounded=True, radius=0.08, shadow=True)
    rect(s, x, 2.25, cw2, 0.12, fill=C["coral"])
    txt(s, x+0.18, 2.6, cw2-0.36, 0.9, [P([R(big, 33, C["coral"], True)], sa=0)], align=PP_ALIGN.CENTER)
    txt(s, x+0.18, 3.55, cw2-0.36, 0.7, [P([R(lab, 13, C["ink"], True)], sa=0)], align=PP_ALIGN.CENTER)
    txt(s, x+0.18, 4.25, cw2-0.36, 0.45, [P([R(sub, 10, C["slate"])], sa=0)], align=PP_ALIGN.CENTER)
rect(s, MARGIN, 5.25, CW, 1.1, fill=C["dusk"], rounded=True, radius=0.06)
txt(s, MARGIN+0.4, 5.45, CW-0.8, 0.8, [
    P([R("취업·경쟁·막연한 미래 — 외재적 압박은 개인이 바꾸기 어렵다. ", 15, C["white"]),
       R("그래서 우리는 ‘바꿀 수 있는’ 내재적 요인에 주목했다.", 15, C["amber"], True)], sa=0, line=1.2)],
    anchor=MSO_ANCHOR.MIDDLE)
footer(s)
note(s, "먼저 문제의 규모입니다. 우울증 진료 1위가 20대이고 5년 새 90% 늘었습니다. 번아웃·정신건강 문제 경험률도 "
        "가파르게 오르는데, 정작 서비스 이용률은 16%에 그칩니다. 취업·경쟁 같은 외재적 압박은 우리가 바꾸기 어렵습니다. "
        "그래서 바꿀 수 있는 내재적 요인에 주목했습니다.")
record(2, "문제정의(1) 청년 마음건강 급증",
       ["우울증 20대 19.4%(’18→’22 +90%)", "번아웃 25–29세 34.8%", "정신건강 경험률 73.6%", "서비스 이용률 16.2%"],
       "외재적 압박은 못 바꾼다 → 내재적 요인(감정인지)로 좁힌다.")

# ---------------------------------------------------------------------
# SLIDE 3 — 문제정의 (2) 감정인지 결여
# ---------------------------------------------------------------------
s = slide()
bg(s, C["parch"])
kicker_title(s, "Ⅰ. 문제정의 — IDEA", "진짜 원인은 ‘감정인지의 결여’다", "03")
# 좌: 논리 체인
chain = [
    ("막연하게 힘들지만", "‘기분이 나쁘다’ 이상으로 구분하지 못함"),
    ("감정을 인지·언어화 못함", "감정표현불능 ↔ 우울·불안과 유의한 상관"),
    ("자기효능감 저하", "→ 무기력·불안의 악순환 (선행 요인)"),
]
for i, (a, b) in enumerate(chain):
    y = 2.2 + i * 1.0
    rect(s, MARGIN, y, 5.7, 0.86, fill=C["white"], rounded=True, radius=0.1, shadow=True)
    rect(s, MARGIN, y, 0.5, 0.86, fill=C["coral"], rounded=False)
    txt(s, MARGIN+0.12, y, 0.5, 0.86, [P([R(str(i+1), 18, C["white"], True)], sa=0)],
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    txt(s, MARGIN+0.7, y+0.13, 4.9, 0.7, [
        P([R(a, 15, C["ink"], True)], sa=1), P([R(b, 11, C["slate"])], sa=0)])
    if i < 2:
        txt(s, MARGIN+2.6, y+0.84, 0.6, 0.2, [P([R("▼", 11, C["amber"], True)], sa=0)])
# 우: 사전설문 + 인용
rect(s, 7.2, 2.2, CW-6.5, 2.05, fill=C["dusk"], rounded=True, radius=0.07)
txt(s, 7.5, 2.36, CW-7.0, 1.85, [
    P([R("사전 설문  ·  n=22 · 2030 청년", 12, C["amber"], True)], sa=6),
    P([R("68%", 30, C["white"], True),
       R("  가 감정 언어화에", 14, C["beige"])], sa=2, line=1.1),
    P([R("‘보통 이하(3점 ↓)’ 로 응답", 13, C["beige"])], sa=0, line=1.1)])
rect(s, 7.2, 4.45, CW-6.5, 1.9, fill=C["white"], rounded=True, radius=0.07, shadow=True)
txt(s, 7.5, 4.7, CW-7.1, 1.5, [
    P([R("“", 26, C["coral"], True)], sa=0),
    P([R("힘든 건 아는데, 정확히 뭐 때문인지는 모르겠어요. "
         "남한테 꺼내지도 못하고 혼자 버티게 돼요.”", 14, C["ink"])], sa=4, line=1.3),
    P([R("— 2030 타깃 사용자 인터뷰", 11, C["slate"])], sa=0)])
footer(s)
note(s, "핵심은 감정인지의 결여입니다. 감정을 구분·언어화하지 못하는 ‘감정표현불능’은 우울·불안과 유의한 상관이 있고, "
        "결과가 아니라 선행 요인입니다. 실제로 사전 설문에서 68%가 ‘내 감정을 말로 표현할 수 있다’에 보통 이하로 답했습니다. "
        "인터뷰에서도 ‘힘든 건 아는데 뭐 때문인지 모르겠다’는 말이 반복됐습니다.")
record(3, "문제정의(2) 감정인지 결여",
       ["감정인지=선행요인(감정표현불능↔우울·불안)", "사전설문 68%가 감정 언어화 3점 이하", "인터뷰: ‘힘든 건 아는데 뭐 때문인지 모름’"],
       "감정인지를 핵심 개입점으로 설정한 근거(통계+인용).")

# ---------------------------------------------------------------------
# SLIDE 4 — 개입지점 (왜 기록인가 + 두 장벽)
# ---------------------------------------------------------------------
s = slide()
bg(s, C["parch"])
kicker_title(s, "Ⅰ. 개입지점 — NORTH STAR", "왜 ‘기록’인가, 그리고 무엇이 막는가", "04")
rect(s, MARGIN, 2.2, CW, 1.0, fill=C["mint"], rounded=True, radius=0.06)
txt(s, MARGIN+0.35, 2.2, CW-0.7, 1.0, [
    P([R("발견  ", 13, C["white"], True),
       R("감정 표현적 글쓰기(expressive writing)는 감정 인지·지능을 향상시킨다 ", 15, C["white"]),
       R("(Pennebaker & Beall 1986; Smyth 1998)", 12, C["dawn"])], sa=0, line=1.2)],
    anchor=MSO_ANCHOR.MIDDLE)
txt(s, MARGIN, 3.45, CW, 0.4, [P([R("그러나 청년은 기록을 지속하지 않는다 — 두 가지 장벽", 15, C["ink"], True)], sa=0)])
barr = [("① 기록이 귀찮다", "앱 열고·타이핑하고·저장하는 마찰. mood-tracking 앱 평균 이탈률 28%, 알림은 오히려 역효과 (JMIR 2026, N=17,123)"),
        ("② 기록이 뭘 돌려주는지 모른다", "기록이 어디 쌓이고 내가 어떻게 달라지는지 체감되지 않는다")]
for i, (t, d) in enumerate(barr):
    x = MARGIN + i * (CW/2 + 0.15)
    rect(s, x, 3.95, CW/2-0.15, 1.3, fill=C["white"], rounded=True, radius=0.08, shadow=True)
    txt(s, x+0.3, 4.15, CW/2-0.75, 1.0, [
        P([R(t, 16, C["coral"], True)], sa=4), P([R(d, 12, C["slate"])], sa=0, line=1.25)])
rect(s, MARGIN, 5.5, CW, 0.95, fill=C["dusk"], rounded=True, radius=0.06)
txt(s, MARGIN+0.35, 5.5, CW-0.7, 0.95, [
    P([R("북극성(North Star)  ", 13, C["amber"], True),
       R("일상 기록이 자연스럽게 ‘감정 패턴 자기인지’로 이어지는 작은 성공을 반복하게 한다", 15, C["white"], True)],
      sa=0, line=1.2)], anchor=MSO_ANCHOR.MIDDLE)
footer(s)
note(s, "감정 표현적 글쓰기는 감정 인지를 키운다는 연구가 있습니다. 그래서 ‘기록’을 개입 수단으로 삼되, 두 장벽을 풀어야 했습니다. "
        "하나는 귀찮음 — 일반 무드트래커는 28%가 이탈하고 알림은 오히려 역효과입니다. 다른 하나는 ‘기록이 뭘 돌려주는지 모름’입니다. "
        "그래서 북극성을 ‘기록이 감정 패턴 인지로 이어지는 작은 성공’으로 잡았습니다.")
record(4, "개입지점 — 왜 기록인가 + 두 장벽",
       ["근거: expressive writing(Pennebaker·Smyth)", "장벽① 귀찮음(무드앱 이탈 28%, 알림 역효과)", "장벽② 기록이 뭘 돌려주는지 모름", "북극성: 기록→감정패턴 자기인지"],
       "두 장벽을 동시에 풀어야 한다는 설계 출발점.")

# ---------------------------------------------------------------------
# SLIDE 5 — 솔루션 컨셉
# ---------------------------------------------------------------------
s = slide()
bg(s, C["dawn"])
kicker_title(s, "Ⅱ. 솔루션 — PROTOTYPE", "카카오톡으로 ‘채집’하고, 웹에서 ‘돌아본다’", "05")
flow = [("카카오톡 채집", "한 줄·사진이면 끝.\n별도 앱 설치 X", C["amber"]),
        ("AI 감정 원석화", "기록을 감정 원석으로\n분류·시각화", C["coral"]),
        ("웹 회고", "도감·캘린더·세공소에서\n내 감정 패턴 회고", C["mint"])]
for i, (t, d, col) in enumerate(flow):
    x = MARGIN + i * 2.95
    rect(s, x, 2.35, 2.6, 1.65, fill=C["white"], rounded=True, radius=0.1, shadow=True)
    rect(s, x, 2.35, 2.6, 0.12, fill=col)
    txt(s, x+0.2, 2.62, 2.2, 0.5, [P([R(t, 16, C["ink"], True)], sa=0)], align=PP_ALIGN.CENTER)
    for j, ln in enumerate(d.split("\n")):
        txt(s, x+0.2, 3.08+j*0.32, 2.2, 0.35, [P([R(ln, 11.5, C["slate"])], sa=0)], align=PP_ALIGN.CENTER)
    if i < 2:
        txt(s, x+2.55, 2.95, 0.45, 0.5, [P([R("→", 22, C["coral"], True)], sa=0)], align=PP_ALIGN.CENTER)
txt(s, MARGIN, 4.25, 5.4, 1.6, [
    P([R("핵심 원칙", 13, C["coral"], True)], sa=6),
    P([R("가벼운 기록", 16, C["ink"], True), R("  ·  ", 14, C["muted"]),
       R("쉬운 접근성", 16, C["ink"], True), R("  ·  ", 14, C["muted"]),
       R("빠른 자기인지", 16, C["ink"], True)], sa=10),
    P([R("기록을 ‘노동’이 아니라 ‘채집’이라는 가벼운 놀이로. ", 13, C["slate"]),
       R("모을수록 도감이 채워지고, 감정은 ‘다룰 수 있는 것’이 된다.", 13, C["slate"])], sa=0, line=1.3)])
pic_fit(s, asset("ai", "chatbot", "gems", "all_gems.png"), 6.5, 4.2, CW-5.8, 2.5, ha="center", va="top")
txt(s, 6.5, 6.7, CW-5.8, 0.3, [P([R("▲ 감정 도감 — 10감정 × 4등급 원석", 10, C["slate"])], sa=0)], align=PP_ALIGN.CENTER)
footer(s)
note(s, "솔루션은 세 단계입니다. 카카오톡에 한 줄이나 사진을 보내면 채집이 끝나고, AI가 감정 원석으로 바꿔주며, 웹에서 도감·캘린더로 "
        "돌아봅니다. 원칙은 가벼운 기록·쉬운 접근·빠른 자기인지. 기록을 노동이 아니라 ‘채집’이라는 놀이로 바꿨고, 모을수록 도감이 채워집니다.")
record(5, "솔루션 컨셉",
       ["채집(카톡)→원석화(AI)→회고(웹)", "원칙: 가벼운 기록·쉬운 접근·빠른 자기인지", "도감: 10감정×4등급 게이미피케이션"],
       "기록을 채집 놀이로. 도감 이미지로 게이미피케이션 시각화(피드백 반영).")

# ---------------------------------------------------------------------
# SLIDE 6 — 시스템 / 프로토타입
# ---------------------------------------------------------------------
s = slide()
bg(s, C["parch"])
kicker_title(s, "Ⅱ. 프로토타입 — SYSTEM", "프로토타입 구조와 기록 모드", "06")
# 좌: 아키텍처 다이어그램 (도형)
ax0, ay0 = MARGIN, 2.2
nodes = [("카카오톡 채널\n챗봇 ‘닥토’", ax0, ay0, C["amber"]),
         ("백엔드 + AI 분류\n(25감정→5계열→10코드)", ax0, ay0+1.25, C["coral"]),
         ("웹 (PWA)\n도감·세공소·캘린더·감정분석", ax0, ay0+2.5, C["mint"])]
for t, x, y, col in nodes:
    rect(s, x, y, 5.2, 1.0, fill=C["white"], rounded=True, radius=0.1, shadow=True)
    rect(s, x, y, 0.16, 1.0, fill=col)
    lines = t.split("\n")
    txt(s, x+0.35, y+0.16, 4.7, 0.8, [
        P([R(lines[0], 14.5, C["ink"], True)], sa=1),
        P([R(lines[1], 11, C["slate"])], sa=0)])
for yy in (ay0+1.0, ay0+2.25):
    txt(s, ax0+2.4, yy+0.02, 0.6, 0.25, [P([R("↕", 13, C["coral"], True)], sa=0)], align=PP_ALIGN.CENTER)
rect(s, ax0, ay0+3.75, 5.2, 0.7, fill=C["dusk"], rounded=True, radius=0.1)
txt(s, ax0+0.3, ay0+3.75, 4.6, 0.7, [
    P([R("Wizard-of-Oz 5일 MVP — ", 12, C["amber"], True),
       R("AI 리소스 절감 + 유저 반응 정밀 관찰", 12, C["white"])], sa=0, line=1.15)],
    anchor=MSO_ANCHOR.MIDDLE)
# 우: 모드 이분화 + 와이어프레임
rect(s, 6.9, 2.2, 3.05, 4.25, fill=C["dusk"], rounded=True, radius=0.05)
pic_fit(s, asset("design", "wireframes", "W-02-home-field-dusk.png"), 7.02, 2.32, 2.8, 4.0, ha="center", va="middle")
txt(s, 6.9, 6.5, 3.05, 0.3, [P([R("▲ 웹 홈 ‘마음 산책길’", 10, C["slate"])], sa=0)], align=PP_ALIGN.CENTER)
rect(s, 10.15, 2.2, CW-9.45, 4.25, fill=C["white"], rounded=True, radius=0.06, shadow=True)
txt(s, 10.4, 2.4, CW-9.95, 4.0, [
    P([R("기록 모드 이분화", 14, C["coral"], True)], sa=8),
    P([R("감정분류 모드", 13, C["ink"], True)], sa=1),
    P([R("챗봇과 함께 감정을 분류·확정·저장", 11, C["slate"])], sa=10, line=1.2),
    P([R("단순기록 모드", 13, C["ink"], True)], sa=1),
    P([R("응답 최소화, 백그라운드 분류", 11, C["slate"])], sa=10, line=1.2),
    P([R("→ 잦은 응답 요구로 인한 ", 11, C["slate"]),
       R("이탈을 줄이는 설계", 11, C["coral"], True)], sa=0, line=1.2)])
footer(s)
note(s, "구조는 카카오톡 챗봇, 백엔드 AI 분류, 웹 세 층입니다. 감정은 25종을 5계열·10코드로 매핑합니다. MVP는 개발 리소스를 아끼려 "
        "운영자가 직접 분류·응대하는 Wizard-of-Oz 방식으로 5일 운영했습니다. 또 사용자 피로를 줄이려 ‘감정분류/단순기록’ 두 모드로 나눴습니다.")
record(6, "시스템/프로토타입",
       ["아키텍처: 챗봇·백엔드AI·웹", "25감정→5계열→10코드", "WoZ 5일 MVP", "모드 이분화(감정분류/단순기록)로 이탈↓"],
       "플로우·모드 설명(피드백 반영). 실제 홈 화면 와이어프레임.")

# ---------------------------------------------------------------------
# SLIDE 7 — 핵심 기능 × HCI 설계 근거  (교수 피드백 정조준)
# ---------------------------------------------------------------------
s = slide()
bg(s, C["parch"])
kicker_title(s, "Ⅱ. 기능 × HCI 설계 근거", "‘왜 이렇게 만들었나’ — 근거 있는 설계", "07")
cards = [
    ("부정감정 자기인지 질문", "저장 직후 1문항. 게이트 4·트리거 3 조건으로 ‘새로움·지속·재발’ 변곡점에만 개입.",
     "감정 분화(Kashdan 2015) · 감정 명명→편도체 하향(Lieberman 2007)"),
    ("홈: 원형 호수 + 조이스틱", "선형 ‘길’을 원형 배치로, 터치 대신 조이스틱 이동/근접 질문 클릭으로 분리.",
     "이동·선택 모드 혼동 회피 · fat-finger · 48dp 터치타깃 (Fitts·Material)"),
    ("캘린더·Recap·감정분석", "순간 인지 → 누적 → 패턴 회고로 이어지는 3층 회고 구조.",
     "긍정 세이버링(Bryant & Veroff) · 반추 회피형 질문(Nolen-Hoeksema)"),
]
cw3 = (CW - 0.6) / 3
for i, (t, d, basis) in enumerate(cards):
    x = MARGIN + i * (cw3 + 0.3)
    rect(s, x, 2.25, cw3, 3.95, fill=C["white"], rounded=True, radius=0.07, shadow=True)
    rect(s, x, 2.25, cw3, 0.7, fill=[C["coral"], C["mint"], C["amber"]][i], rounded=True, radius=0.07)
    rect(s, x, 2.6, cw3, 0.35, fill=[C["coral"], C["mint"], C["amber"]][i])
    txt(s, x+0.22, 2.36, cw3-0.44, 0.6, [P([R(t, 14.5, C["white"], True)], sa=0)], anchor=MSO_ANCHOR.MIDDLE)
    txt(s, x+0.22, 3.15, cw3-0.44, 1.7, [P([R(d, 12.5, C["ink"])], sa=0, line=1.3)])
    rect(s, x+0.22, 5.0, cw3-0.44, 1.0, fill=C["dawn"], rounded=True, radius=0.1)
    txt(s, x+0.4, 5.1, cw3-0.8, 0.85, [
        P([R("HCI 근거", 10, C["coral"], True)], sa=2),
        P([R(basis, 10.5, C["slate"])], sa=0, line=1.2)])
footer(s)
note(s, "교수님께서 ‘기능 설명 같다’고 하셔서, 이번엔 ‘왜 이렇게 만들었나’를 HCI 근거로 보여드립니다. 자기인지 질문은 감정 분화 이론과 "
        "감정 명명의 편도체 하향 효과에 기반하고, 변곡점에만 개입하도록 설계했습니다. 홈의 원형 호수와 조이스틱은 이동·선택 모드 혼동과 "
        "fat-finger 문제를 피하려 Fitts의 법칙과 48dp 터치타깃 가이드를 따랐습니다. 회고 구조도 세이버링·반추 회피 이론에 근거합니다.")
record(7, "핵심 기능 × HCI 설계 근거",
       ["자기인지 질문: Kashdan·Lieberman", "원형 호수+조이스틱: Fitts·fat-finger·48dp", "캘린더/Recap/감정분석: 세이버링·반추회피"],
       "★교수 ‘HCI 관점’ 정조준. 각 설계 결정에 학술 근거.")

# ---------------------------------------------------------------------
# SLIDE 8 — 실험: 연구질문 & 가설
# ---------------------------------------------------------------------
s = slide()
bg(s, C["dusk"], C["ink"], angle=120)
rect(s, MARGIN, 0.64, 0.12, 0.62, fill=C["amber"])
txt(s, MARGIN+0.28, 0.62, CW-0.28, 0.4, [P([R("Ⅲ. 실험 설계 — EXPERIMENT DESIGN", 13, C["amber"], True)], sa=0)])
txt(s, MARGIN+0.28, 0.96, CW-0.28, 0.7, [P([R("연구 질문과 가설", 27, C["white"], True)], sa=0)])
rect(s, MARGIN, 2.1, CW, 1.15, fill=C["coral"], rounded=True, radius=0.06)
txt(s, MARGIN+0.4, 2.1, CW-0.8, 1.15, [
    P([R("RQ  ", 14, C["dawn"], True),
       R("카카오톡 기반의 낮은 진입장벽 기록 + 원석·Recap 회고 경험이 ", 15, C["white"]),
       R("사용자의 감정 자기인지", 15, C["amber"], True),
       R(" 를 유도할 수 있는가?", 15, C["white"])], sa=0, line=1.3)],
    anchor=MSO_ANCHOR.MIDDLE)
rect(s, MARGIN, 3.5, CW, 1.0, fill=C["white"], rounded=True, radius=0.06, shadow=True)
txt(s, MARGIN+0.4, 3.5, CW-0.8, 1.0, [
    P([R("H1  ", 14, C["coral"], True),
       R("감정인지 능력이 향상되면 막연한 불안감이 줄고 자기효능감이 높아진다.", 15, C["ink"], True)], sa=2, line=1.2),
    P([R("근거: 감정 명명→편도체 하향(Lieberman 2007) · 감정 분화(Kashdan 2015) · 정서인식→자기효능감(이한우·이미희 2013)",
         11, C["slate"])], sa=0, line=1.2)],
    anchor=MSO_ANCHOR.MIDDLE)
# 변수 박스
iv = ["서비스 사용 전 → 후 (within-subject)", "1차 MVP → 2차 MVP (분류 로직 개선)"]
dv = ["감정인지 자기보고(Kashdan) · 정서 K-PANAS", "SUS · AI 신뢰도 · 전환율 · 리텐션 · NPS"]
for i, (head, arr, col) in enumerate([("독립변수 (IV)", iv, C["amber"]), ("종속변수 (DV)", dv, C["mint"])]):
    x = MARGIN + i * (CW/2 + 0.15)
    rect(s, x, 4.75, CW/2-0.15, 1.7, fill=rgb("#2E3149"), rounded=True, radius=0.06)
    txt(s, x+0.3, 4.92, CW/2-0.75, 0.4, [P([R(head, 13, col, True)], sa=0)])
    bullets(s, x+0.3, 5.35, CW/2-0.75, 1.0, arr, size=12, color=C["beige"], gap=6, markcolor=col)
note(s, "실험의 연구 질문은 ‘저진입 기록과 회고 경험이 감정 자기인지를 유도하는가’입니다. 가설 H1은 ‘감정인지가 오르면 막연한 불안이 줄고 "
        "자기효능감이 오른다’이고, 세 편의 연구에 근거합니다. 독립변수는 ‘사용 전후’와 ‘1차→2차 개선’, 종속변수는 감정인지·SUS·전환·리텐션·NPS입니다.")
record(8, "실험: 연구질문 & 가설",
       ["RQ: 저진입 기록+회고 → 감정 자기인지?", "H1: 감정인지↑→불안↓·자기효능감↑", "IV: 사용 전후 / 1차→2차", "DV: 감정인지·SUS·전환·리텐션·NPS"],
       "가설을 학술 근거와 함께 명시. IV/DV 정식화(교수 피드백).")

# ---------------------------------------------------------------------
# SLIDE 9 — 실험: 방법 (혼합설계)
# ---------------------------------------------------------------------
s = slide()
bg(s, C["parch"])
kicker_title(s, "Ⅲ. 방법 — MIXED METHODS", "참가자 · 측정도구 · KPI", "09")
# 참가자
rect(s, MARGIN, 2.2, 3.7, 2.0, fill=C["white"], rounded=True, radius=0.08, shadow=True)
txt(s, MARGIN+0.25, 2.36, 3.25, 1.9, [
    P([R("참가자", 13, C["coral"], True)], sa=4),
    P([R("2030 청년 (남40/여60)", 12.5, C["ink"], True)], sa=4, line=1.12),
    P([R("사전 22 · 사후 16 (중복제거 14)", 11.5, C["ink"])], sa=4, line=1.12),
    P([R("사전·사후 매칭 ", 11.5, C["ink"]), R("8명 (36%)", 11.5, C["coral"], True)], sa=4, line=1.12),
    P([R("05.18–22 · Wizard-of-Oz", 10.5, C["slate"])], sa=0)])
# 측정도구
rect(s, MARGIN+3.95, 2.2, CW-3.95, 2.0, fill=C["white"], rounded=True, radius=0.08, shadow=True)
txt(s, MARGIN+4.2, 2.38, CW-4.45, 0.4, [P([R("측정 도구 (검증된 척도)", 13, C["mint"], True)], sa=0)])
tools = [("감정인지 5문항", "Kashdan 2015"), ("정서 PA/NA 7문항", "K-PANAS"),
         ("사용성 10문항", "SUS · Brooke 1996"), ("AI 신뢰도 2 · NPS 1", "Likert / 0–10")]
for i, (a, b) in enumerate(tools):
    x = MARGIN+4.2 + (i % 2)*(CW-4.45)/2
    y = 2.82 + (i//2)*0.62
    rect(s, x, y, (CW-4.45)/2-0.2, 0.5, fill=C["dawn"], rounded=True, radius=0.12)
    txt(s, x+0.15, y, (CW-4.45)/2-0.5, 0.5, [
        P([R(a+"  ", 12, C["ink"], True), R(b, 10, C["slate"])], sa=0)], anchor=MSO_ANCHOR.MIDDLE)
# KPI 스트립
rect(s, MARGIN, 4.4, CW, 2.0, fill=C["dusk"], rounded=True, radius=0.05)
txt(s, MARGIN+0.3, 4.55, CW-0.6, 0.4, [P([R("KPI 10개 — 효과 · 사용성 · 지속성을 동시에 측정", 13, C["amber"], True)], sa=0)])
kpis = [("K1 감정인지", "사후 ≥ 사전+0.5"), ("K5 AI 만족도", "≥ 3.5/5"),
        ("K6 AI 수정률", "≤ 30%"), ("K9 카톡→웹 전환", "≥ 50%"),
        ("K3 3회+ 기록", "≥ 20%"), ("K10 오류 이탈", "≤ 20%")]
for i, (a, b) in enumerate(kpis):
    x = MARGIN+0.3 + (i % 3)*((CW-0.6)/3)
    y = 5.05 + (i//3)*0.62
    rect(s, x, y, (CW-0.6)/3-0.2, 0.5, fill=rgb("#2E3149"), rounded=True, radius=0.1)
    txt(s, x+0.18, y, (CW-0.6)/3-0.5, 0.5, [
        P([R(a+"  ", 11.5, C["white"], True), R(b, 10.5, C["mint"])], sa=0)], anchor=MSO_ANCHOR.MIDDLE)
txt(s, MARGIN, 6.55, CW, 0.35, [P([R("데이터 = 사전/사후 설문(Google Form) + 행동 로그(DB) + 사후 인터뷰  ·  삼각측정",
                                     11, C["slate"], True)], sa=0)])
footer(s)
note(s, "방법은 혼합설계입니다. 2030 청년을 모집해 사전 22명, 사후 16명, 사전·사후 매칭은 8명이었습니다. 척도는 모두 검증된 것 — "
        "감정인지는 Kashdan, 정서는 K-PANAS, 사용성은 SUS를 썼습니다. KPI 10개로 효과·사용성·지속성을 함께 봤고, 설문·행동로그·인터뷰를 "
        "삼각측정했습니다. 표본이 작다는 한계는 뒤에서 정직하게 다룹니다.")
record(9, "실험: 방법(혼합설계)",
       ["참가자 사전22/사후16/매칭8(36%)", "척도: Kashdan·K-PANAS·SUS·NPS", "KPI 10개", "설문+행동로그+인터뷰 삼각측정"],
       "검증된 척도와 KPI, 삼각측정. 표본 한계는 뒤에서.")

# ---------------------------------------------------------------------
# SLIDE 10 — Before/After & A/B 설계
# ---------------------------------------------------------------------
s = slide()
bg(s, C["parch"])
kicker_title(s, "Ⅲ. 비교 설계 — BEFORE/AFTER · A/B", "두 가지 비교 실험을 설계했다", "10")
rect(s, MARGIN, 2.25, CW/2-0.2, 3.9, fill=C["white"], rounded=True, radius=0.07, shadow=True)
rect(s, MARGIN, 2.25, CW/2-0.2, 0.6, fill=C["coral"], rounded=True, radius=0.07)
rect(s, MARGIN, 2.55, CW/2-0.2, 0.3, fill=C["coral"])
txt(s, MARGIN+0.25, 2.3, CW/2-0.7, 0.55, [P([R("(a) Before / After — 분류 로직 개선", 14, C["white"], True)], sa=0)],
    anchor=MSO_ANCHOR.MIDDLE)
bullets(s, MARGIN+0.3, 3.1, CW/2-0.8, 2.9, [
    ("조작: 1차 → 2차 분류 로직", "supervisor 2차 검증 · 멀티감정 · 유해키워드 FP 축소 · 명령어 정규화"),
    ("측정: 같은 잣대(사용자 정정)", "web_reviewed 행의 AI=확정 일치율"),
    ("같은 사용자·같은 채널 — 시점만 분리", "운영 환경 그대로의 자연 실험"),
], size=12.5, gap=8)
rect(s, MARGIN+CW/2+0.05, 2.25, CW/2-0.2, 3.9, fill=C["white"], rounded=True, radius=0.07, shadow=True)
rect(s, MARGIN+CW/2+0.05, 2.25, CW/2-0.2, 0.6, fill=C["mint"], rounded=True, radius=0.07)
rect(s, MARGIN+CW/2+0.05, 2.55, CW/2-0.2, 0.3, fill=C["mint"])
txt(s, MARGIN+CW/2+0.3, 2.3, CW/2-0.7, 0.55, [P([R("(b) A/B — 알림(리마인드) 실험", 14, C["white"], True)], sa=0)],
    anchor=MSO_ANCHOR.MIDDLE)
bullets(s, MARGIN+CW/2+0.35, 3.1, CW/2-0.8, 2.9, [
    ("조건 A: 행동 요청형", "“오늘 기록해주세요”"),
    ("조건 B: 회고 가치형", "“오늘 한 줄이면 주간 감정 흐름이 완성돼요”"),
    ("지표: 재방문·기록률 (3일 측정)", "멘토/교수 권고를 실험으로 정식화"),
], size=12.5, gap=8)
footer(s)
note(s, "비교는 두 가지로 설계했습니다. 첫째 Before/After — 1차에서 2차로 분류 로직을 개선하고 같은 사용자·같은 채널에서 정정 기반 동의율을 "
        "비교한 자연 실험입니다. 둘째는 알림 A/B입니다. 교수님과 멘토님이 ‘알림 A/B를 3일 해보라’ 하셨는데, 행동 요청형 대 회고 가치형 문구를 "
        "비교하도록 정식화했습니다. 다음 장부터 실제 결과입니다.")
record(10, "Before/After & A/B 설계",
       ["(a) 1차→2차 분류 개선 자연실험", "(b) 알림 A/B: 행동요청형 vs 회고가치형(3일)", "교수/멘토 ‘A/B 해보라’ 정식화"],
       "교수 피드백(A/B)을 실험으로 명문화.")

# ---------------------------------------------------------------------
# Results section divider
# ---------------------------------------------------------------------
def chart_slide(tnum, kicker, title, chart, takeaway, sub_left=None, extra=None):
    s = slide(); bg(s, C["parch"])
    kicker_title(s, kicker, title, tnum)
    pic_fit(s, os.path.join(CHARTS, chart), 6.7, 2.2, CW-6.0, 4.0, ha="center", va="middle")
    if sub_left:
        bullets(s, MARGIN, 2.35, 5.6, 3.0, sub_left, size=13.5, gap=10)
    if extra:
        rect(s, MARGIN, 5.55, 5.8, 0.95, fill=C["dusk"], rounded=True, radius=0.06)
        txt(s, MARGIN+0.3, 5.55, 5.4, 0.95, [P([R(extra, 12.5, C["white"])], sa=0, line=1.25)],
            anchor=MSO_ANCHOR.MIDDLE)
    rect(s, 6.7, 6.45, CW-6.0, 0.0)  # spacer
    txt(s, 6.7, 6.5, CW-6.0, 0.4, [P([R("▲ " + takeaway, 11, C["slate"], True)], sa=0)], align=PP_ALIGN.CENTER)
    footer(s)
    return s

# SLIDE 11 — 결과① 감정인지
s = chart_slide("11", "Ⅳ. 결과 — USER STUDY", "결과① 감정인지: 가설의 방향을 확인", "emotion_recognition.png",
    "사전·사후 매칭 8명 / 효과크기 d=0.81",
    sub_left=[
        ("사전 3.60 → 사후 3.90 (+0.30)", "감정인지 자기보고 5점 척도 (Kashdan)"),
        ("8명 중 6명 상승 · 1명 유지 · 1명 하락", "초기 감정인지가 낮았던 사람일수록 더 크게 상승"),
        ("Cohen’s d = 0.81 (큰 효과크기)", "t-test p=0.056 — 통계적 유의 직전"),
        ("정직하게: 탐색적 결과", "표본이 작아 ‘입증’이 아닌 ‘방향성 확인’"),
    ],
    extra="가설 H1의 방향(감정인지↑)이 데이터에서 관찰됐다. 효과크기는 크지만, 유의성은 표본 확대가 필요하다.")
note(s, "첫 결과입니다. 감정인지 자기보고가 사전 3.60에서 사후 3.90으로 올랐고, 8명 중 6명이 상승했습니다. 효과크기 d는 0.81로 크지만 "
        "p값은 0.056으로 유의 직전입니다. 표본이 작아 ‘입증’이라기보다 가설의 ‘방향’을 확인한 탐색적 결과로 정직하게 말씀드립니다.")
record(11, "결과① 감정인지", ["3.60→3.90(+0.30), 6/8 상승", "Cohen’s d=0.81, p=0.056", "탐색적·방향성 확인"],
       "핵심 가설 방향 확인. 효과크기 크나 유의 직전 — 정직하게.")

# SLIDE 12 — 결과② 정확도 Before/After  (FINDING 포맷)
s = slide(); bg(s, C["dusk"], C["ink"], angle=120)
rect(s, MARGIN, 0.64, 0.12, 0.62, fill=C["mint"])
txt(s, MARGIN+0.28, 0.62, CW, 0.4, [P([R("Ⅳ. 결과 — FINDING", 13, C["mint"], True)], sa=0)])
txt(s, MARGIN+0.28, 0.96, CW, 0.7, [P([R("결과② AI 분류 정확도 — 모델은 문제가 아니다", 26, C["white"], True)], sa=0)])
# 큰 숫자
big = [("97.3%", "2차 사용자 동의율", C["mint"], 34), ("2.7%", "AI 수정률 (목표 ≤30%)", C["amber"], 34),
       ("41.7 → 97.3", "개선 직전 → 2차 (%)", C["coral"], 26)]
for i, (b, l, col, bs) in enumerate(big):
    x = MARGIN + i*(CW/3)
    txt(s, x, 2.15, CW/3-0.25, 0.95, [P([R(b, bs, col, True)], sa=0)], anchor=MSO_ANCHOR.BOTTOM)
    txt(s, x, 3.15, CW/3-0.3, 0.5, [P([R(l, 12.5, C["beige"])], sa=0)])
pic_fit(s, os.path.join(CHARTS, "accuracy_beforeafter.png"), MARGIN, 3.75, 5.7, 2.7, ha="left", va="top")
rect(s, 6.9, 3.8, CW-6.2, 2.55, fill=rgb("#2E3149"), rounded=True, radius=0.06)
txt(s, 7.15, 3.98, CW-6.7, 2.3, [
    P([R("채널 96.6%  ·  웹 재검수 83.7%  →  실질 84–97%", 13, C["white"], True)], sa=8, line=1.2),
    P([R("개선 직전(05.13–20) 정정 다발 → 2차 배포 후 정정 0건", 12, C["beige"])], sa=8, line=1.2),
    P([R("남은 약점: ", 12, C["beige"]), R("‘짜증(annoyance)’ 정밀도 50%", 12, C["coral"], True),
       R(" — 나머지 ≥78%", 12, C["beige"])], sa=8, line=1.2),
    P([R("→ 모델 품질이 아니라 ‘다음 단계 전환’이 진짜 병목", 12.5, C["amber"], True)], sa=0, line=1.2)])
note(s, "둘째, AI 분류 정확도입니다. 2차 사용자 동의율이 97.3%, 수정률은 2.7%로 목표(30% 이하)를 10배 이상 초과했습니다. 개선 직전 41.7%에서 "
        "2차 97.3%로 뛰었고 배포 후 정정은 0건이었습니다. 약점은 ‘짜증’ 클래스 하나로 정밀도 50%였습니다. 결론은 — 모델 품질은 이미 충분하고, "
        "진짜 병목은 다음 단계 전환이라는 점입니다.")
record(12, "결과② 정확도 Before/After",
       ["2차 동의율 97.3%, 수정률 2.7%(목표≤30%)", "41.7%→97.3%, 배포 후 정정 0", "짜증 정밀도 50% 약점", "병목은 모델이 아니라 전환"],
       "FINDING 포맷 큰 숫자. 모델은 문제 아님을 선언.")

# SLIDE 13 — 결과③ 사용성 · 회고 · 전환
s = slide(); bg(s, C["parch"])
kicker_title(s, "Ⅳ. 결과 — USER STUDY", "결과③ 사용자는 ‘분석’보다 ‘회고’를 더 높이 평가했다", "13")
pic_fit(s, os.path.join(CHARTS, "sus.png"), MARGIN, 2.2, 5.2, 2.85, ha="left", va="top")
txt(s, MARGIN, 5.06, 5.2, 0.45, [P([R("SUS 71.25 — 업계 평균(68) 상회 (편차 40–100)", 11.5, C["slate"], True)], sa=0, line=1.15)])
pic_fit(s, os.path.join(CHARTS, "recap_features.png"), 6.5, 2.2, CW-5.8, 2.85, ha="center", va="top")
txt(s, 6.5, 5.06, CW-5.8, 0.45, [P([R("회고 기능 만족도 — 캘린더(3.77) > AI 분류 만족도(3.36)", 11.5, C["slate"], True)], sa=0, line=1.15)])
rect(s, MARGIN, 5.65, CW, 0.95, fill=C["dusk"], rounded=True, radius=0.06)
txt(s, MARGIN+0.35, 5.65, 2.7, 0.95, [
    P([R("카톡 → 웹 전환", 11, C["amber"], True)], sa=1),
    P([R("설문 93%", 15, C["white"], True), R("  vs  ", 11, C["beige"]), R("행동 23%", 15, C["coral"], True)], sa=0)],
    anchor=MSO_ANCHOR.MIDDLE)
txt(s, MARGIN+4.0, 5.65, CW-4.3, 0.95, [
    P([R("자기보고와 행동의 간극 — ‘들어갔다’는 기억과 실제 로그가 다르다. ", 12.5, C["white"]),
       R("측정 방법이 결론을 바꾼다 (HCI 통찰).", 12.5, C["amber"], True)], sa=0, line=1.2)],
    anchor=MSO_ANCHOR.MIDDLE)
footer(s)
note(s, "셋째, 사용성은 SUS 71.25로 업계 평균 68을 넘었습니다(편차는 큼). 더 중요한 건 회고 기능 평가입니다. 가장 높은 점수는 캘린더 3.77이고, "
        "AI 분류 만족도는 3.36으로 목표 3.5에 못 미쳤습니다. 즉 사용자는 AI 분석보다 자기 기록을 다시 보는 ‘회고’를 더 높이 평가했습니다. "
        "전환율도 설문에선 93%였지만 행동 로그는 23%로, 자기보고와 행동의 간극이라는 HCI 교훈을 남겼습니다.")
record(13, "결과③ 사용성·회고·전환",
       ["SUS 71.25(>68)", "회고기능: 캘린더 3.77 최고 > AI 만족 3.36", "전환 93%(설문) vs 23%(행동)", "사용자는 분석보다 회고를 원함"],
       "캘린더 최고점이 ‘회고 서비스’ 피벗의 근거. 전환 간극은 측정 방법 교훈.")

# SLIDE 14 — 결과④ 리텐션 & 이탈
s = slide(); bg(s, C["parch"])
kicker_title(s, "Ⅳ. 결과 — 정직한 부정 결과", "결과④ 리텐션과 이탈: 문제는 기능이 아니라 습관", "14")
pic_fit(s, os.path.join(CHARTS, "retention.png"), MARGIN, 2.3, 5.5, 2.9, ha="left", va="top")
pic_fit(s, os.path.join(CHARTS, "nps.png"), MARGIN+0.2, 5.0, 2.2, 2.0, ha="left", va="top")
pic_fit(s, os.path.join(CHARTS, "churn_reasons.png"), 6.9, 2.3, CW-6.2, 2.9, ha="center", va="top")
rect(s, 6.9, 5.3, CW-6.2, 1.15, fill=C["dusk"], rounded=True, radius=0.06)
txt(s, 7.15, 5.3, CW-6.7, 1.15, [
    P([R("활동자 41%가 단 하루. NPS −28.6.", 13, C["amber"], True)], sa=4, line=1.2),
    P([R("이탈 사유 1·2위 = 잊어버림·귀찮음. 사용자는 서비스를 싫어해 떠난 게 아니라, "
         "기록을 떠올리지 못했고 기록할 이유를 못 느꼈다.", 11.5, C["white"])], sa=0, line=1.25)])
txt(s, MARGIN+2.55, 5.2, 3.0, 1.6, [
    P([R("리텐션", 11, C["coral"], True)], sa=3),
    P([R("2일+ 57% · 4일+ 21%", 12, C["ink"], True)], sa=3, line=1.2),
    P([R("그러나 1일 이탈 41%가 최빈값", 11, C["slate"])], sa=0, line=1.2)])
footer(s)
note(s, "넷째, 정직한 부정 결과입니다. 활동 사용자의 41%가 단 하루만 쓰고 떠났고 NPS는 −28.6이었습니다. 이탈 이유 1·2위는 ‘잊어버림’과 "
        "‘귀찮음’이었습니다. 즉 서비스를 싫어해 떠난 게 아니라, 기록을 떠올리지 못했고 기록할 이유를 충분히 못 느낀 겁니다. 문제는 기능이 아니라 습관입니다.")
record(14, "결과④ 리텐션 & 이탈",
       ["활동자 41% 1일 이탈, 2일+ 57%/4일+ 21%", "NPS −28.6", "이탈사유 망각·귀찮음", "문제는 기능이 아니라 습관"],
       "정직한 부정 결과. 습관 형성이 과제.")

# SLIDE 15 — 결과⑤ 정성 인터뷰
s = slide(); bg(s, C["parch"])
kicker_title(s, "Ⅳ. 결과 — 사후 인터뷰", "결과⑤ 사용자는 ‘감정’보다 ‘사건과 맥락’을 보고 싶어 했다", "15")
quotes = [
    ("“감정이 많이 쌓인 날은 ‘무슨 일이 있었지?’ 궁금해서\n캘린더를 가장 자주 봤어요.”", "이찬희 — 캘린더=회고 진입점", C["coral"]),
    ("“밥을 너무 많이 먹어 슬펐을 때… 사실 슬픈 건 아니었는데\nAI한테 여러 번 되물었어요.”", "김인영 — 모호·사건 기록에 약한 AI", C["amber"]),
    ("“짝사랑하던 친구가 데려다준 그 설렘을, 친구에게 말하긴\n구차해서 아보하에 기록하고 싶었어요.”", "황수지 — 사소하지만 의미있는 순간 보관", C["mint"]),
    ("“웹이 있는지도 몰랐어요. 챗봇을 굳이 왜 해야 되는지…\n알림이 와도 그냥 무시했어요.”", "이탈자 — 온보딩·동기 부재", C["slate"]),
]
for i, (q, who, col) in enumerate(quotes):
    x = MARGIN + (i % 2)*(CW/2+0.1)
    y = 2.2 + (i//2)*2.05
    rect(s, x, y, CW/2-0.1, 1.85, fill=C["white"], rounded=True, radius=0.07, shadow=True)
    rect(s, x, y, 0.14, 1.85, fill=col)
    txt(s, x+0.35, y+0.2, CW/2-0.65, 1.2, [P([R(q, 13, C["ink"])], sa=0, line=1.25)])
    txt(s, x+0.35, y+1.42, CW/2-0.65, 0.35, [P([R(who, 11, col, True)], sa=0)])
footer(s)
note(s, "마지막 결과는 인터뷰입니다. 공통점은 사용자가 ‘감정 라벨’보다 ‘그날의 사건과 맥락’을 다시 보고 싶어 한다는 것이었습니다. 이찬희님은 "
        "캘린더를 회고 진입점으로 썼고, 김인영님 사례처럼 AI는 모호하거나 사건만 적은 기록에 약했습니다. 황수지님은 사소하지만 의미 있는 순간을 "
        "보관하고 싶어 했고, 이탈자는 웹의 존재조차 몰랐습니다. 온보딩과 동기 설계가 과제로 드러났습니다.")
record(15, "결과⑤ 정성 인터뷰",
       ["감정<사건·맥락", "AI는 모호·사건 기록에 약함", "사소하지만 의미있는 순간 보관", "온보딩/동기 부재(이탈)"],
       "Verbatim 인용 4개로 정성 통찰.")

# SLIDE 16 — 고찰①
s = slide(); bg(s, C["dusk"], C["ink"], angle=120)
rect(s, MARGIN, 0.64, 0.12, 0.62, fill=C["amber"])
txt(s, MARGIN+0.28, 0.62, CW, 0.4, [P([R("Ⅴ. 고찰 — DISCUSSION", 13, C["amber"], True)], sa=0)])
txt(s, MARGIN+0.28, 0.96, CW, 0.7, [P([R("무엇을 배웠나", 27, C["white"], True)], sa=0)])
ins = [
    ("모델이 아니라 ‘전환·습관’이 병목", "정확도·수정률은 목표를 크게 초과. 잃는 지점은 ‘다음 단계로 안 넘어옴’이었다."),
    ("사람들은 분석받기보다 스스로 알아차리고 싶어 한다", "캘린더(3.77)·Recap이 AI 분석보다 높이 평가됨 → ‘분석 서비스’보다 ‘회고 서비스’."),
    ("자기보고 ↔ 행동의 간극(93% vs 23%)", "측정 방법이 결론을 바꾼다 — 설문만 믿으면 안 된다는 방법론적 교훈."),
    ("→ 다음 방향: 감정+사건을 함께 남기는 ‘회고 서비스’로 피벗", "캘린더 중심 재설계 · 기록 이유를 주는 리마인드 · AI는 단정 대신 후보 제시"),
]
for i, (t, d) in enumerate(ins):
    y = 2.15 + i*1.08
    rect(s, MARGIN, y, CW, 0.95, fill=(C["coral"] if i == 3 else rgb("#2E3149")), rounded=True, radius=0.06)
    txt(s, MARGIN+0.35, y+0.13, CW-0.7, 0.8, [
        P([R(t, 15, C["white"], True)], sa=2, line=1.1),
        P([R(d, 11.5, (C["dawn"] if i == 3 else C["beige"]))], sa=0, line=1.15)])
note(s, "고찰입니다. 첫째, 병목은 모델이 아니라 전환과 습관이었습니다. 둘째, 사람들은 분석받기보다 스스로 알아차리고 싶어 했습니다 — 캘린더와 "
        "Recap이 AI 분석보다 높게 평가됐죠. 셋째, 자기보고와 행동의 간극은 ‘설문만 믿으면 안 된다’는 방법론 교훈입니다. 그래서 우리는 감정과 "
        "사건을 함께 남기는 회고 서비스로 방향을 잡았습니다.")
record(16, "고찰① 무엇을 배웠나",
       ["병목=전환·습관(모델 아님)", "분석<회고(캘린더·Recap 우위)", "자기보고↔행동 간극=방법론 교훈", "→ 회고 서비스로 피벗"],
       "핵심 통찰 4개. 피벗 방향 제시.")

# SLIDE 17 — 고찰② 한계 & 이론
s = slide(); bg(s, C["parch"])
kicker_title(s, "Ⅴ. 고찰 — 한계 & 이론", "한계를 인정하고, 이론으로 설명한다", "17")
rect(s, MARGIN, 2.2, CW/2-0.2, 4.1, fill=C["white"], rounded=True, radius=0.07, shadow=True)
txt(s, MARGIN+0.3, 2.4, CW/2-0.8, 0.4, [P([R("한계 (Limitations)", 14, C["coral"], True)], sa=0)])
bullets(s, MARGIN+0.3, 2.9, CW/2-0.8, 3.3, [
    "작은 표본 · 매칭률 36% → 유의성 확보 못함 (p=0.056)",
    "Wizard-of-Oz 운영 — 완전 자동화 아님",
    "웹 검수 선택편향 · 운영자 트래픽 일부 혼입",
    "1주 테스트 — 장기 습관·리텐션은 미검증",
    "교수 피드백 ‘10명은 적다’ → 표본 확대 필요",
], size=12.5, gap=11)
rect(s, MARGIN+CW/2+0.05, 2.2, CW/2-0.2, 4.1, fill=C["dusk"], rounded=True, radius=0.07)
txt(s, MARGIN+CW/2+0.35, 2.4, CW/2-0.8, 0.4, [P([R("이론적 연결 (Theory)", 14, C["amber"], True)], sa=0)])
bullets(s, MARGIN+CW/2+0.35, 2.9, CW/2-0.85, 3.3, [
    "감정 명명→편도체 하향 (Lieberman 2007) / 감정 분화 (Kashdan 2015)",
    "의인화·자기개방(CASA, Nass) · 베이비 스키마(Glocker 2009)",
    "알림 역효과 (Astill Wright 2026, JMIR, N=17,123)",
    "전문가 자문(토닥토닥): ‘회기 간 공백’ · 예방 중심 · AI는 정보형 대체, 역동형은 인간 영역",
], size=12.5, gap=11, markcolor=C["amber"], color=C["beige"])
footer(s)
note(s, "한계는 정직하게 말씀드립니다. 표본이 작고 매칭률이 36%라 유의성은 확보하지 못했습니다. WoZ였고, 웹 검수 선택편향과 운영자 트래픽도 "
        "일부 있었으며 1주라 장기 리텐션은 검증 못 했습니다. 다만 결과는 이론으로 설명됩니다 — 감정 명명의 편도체 하향, 의인화, 알림 역효과 연구, "
        "그리고 상담 전문가의 ‘회기 간 공백’ 통찰이 우리 설계와 맞닿습니다.")
record(17, "고찰② 한계 & 이론",
       ["한계: 작은 n·매칭36%·WoZ·선택편향·1주", "이론: Lieberman·Kashdan·CASA·JMIR", "전문가 자문: 회기 간 공백·예방 중심"],
       "한계 정직 + 이론 연결로 학술 깊이.")

# SLIDE 18 — 향후 통제 A/B 설계
s = slide(); bg(s, C["parch"])
kicker_title(s, "Ⅵ. 향후 — CONTROLLED EXPERIMENT", "다음 단계: 제대로 된 통제 A/B 실험", "18")
rect(s, MARGIN, 2.2, CW, 0.85, fill=C["mint"], rounded=True, radius=0.06)
txt(s, MARGIN+0.35, 2.2, CW-0.7, 0.85, [
    P([R("H1  ", 13, C["white"], True),
       R("‘회고 가치형’ 개입(캘린더 중심 + 후보 제시형 AI + 회고형 알림)이 ", 13.5, C["white"]),
       R("전환율·리텐션·감정인지", 13.5, C["dawn"], True),
       R(" 를 높인다", 13.5, C["white"])], sa=0, line=1.2)], anchor=MSO_ANCHOR.MIDDLE)
cols3 = [
    ("설계", ["무작위 배정 / 피험자 내 카운터밸런싱", "충분한 표본 · 검정력 분석", "사전등록(pre-registration)"]),
    ("독립변수 (조작)", ["알림: 행동요청형 vs 회고가치형", "AI: 단정형 vs 후보 제시형", "웹: 분석 중심 vs 캘린더 중심"]),
    ("종속변수 (측정)", ["전환율 · 리텐션(행동 로그)", "SUS · NASA-TLX(인지부하)", "감정인지 사전/사후 delta"]),
]
cw3 = (CW-0.6)/3
for i, (h, arr) in enumerate(cols3):
    x = MARGIN + i*(cw3+0.3)
    rect(s, x, 3.25, cw3, 2.4, fill=C["white"], rounded=True, radius=0.07, shadow=True)
    rect(s, x, 3.25, cw3, 0.5, fill=[C["coral"], C["amber"], C["mint"]][i], rounded=True, radius=0.07)
    rect(s, x, 3.5, cw3, 0.25, fill=[C["coral"], C["amber"], C["mint"]][i])
    txt(s, x+0.2, 3.27, cw3-0.4, 0.48, [P([R(h, 13, C["white"], True)], sa=0)], anchor=MSO_ANCHOR.MIDDLE)
    bullets(s, x+0.25, 3.9, cw3-0.5, 1.7, arr, size=11.5, gap=8)
rect(s, MARGIN, 5.85, CW, 0.6, fill=C["dusk"], rounded=True, radius=0.08)
txt(s, MARGIN+0.35, 5.85, CW-0.7, 0.6, [
    P([R("다음 버전 우선순위  ", 12, C["amber"], True),
       R("① 사건+감정 함께 기록  ② 캘린더 중심 회고  ③ AI 맥락 처리  ④ 회고형 리마인드", 12, C["white"])],
      sa=0, line=1.1)], anchor=MSO_ANCHOR.MIDDLE)
footer(s)
note(s, "향후엔 자연 실험을 넘어 제대로 된 통제 실험을 설계했습니다. 무작위 배정과 카운터밸런싱, 검정력 분석, 사전등록을 갖추고, 알림·AI·웹 "
        "세 축을 조작합니다. 종속변수엔 전환·리텐션과 함께 SUS, 인지부하(NASA-TLX), 감정인지 변화량을 둡니다. 다음 버전은 사건+감정 기록, "
        "캘린더 중심 회고, AI 맥락 처리, 회고형 리마인드 순으로 갑니다.")
record(18, "향후 통제 A/B 설계",
       ["무작위배정/카운터밸런싱/검정력/사전등록", "IV: 알림·AI·웹 3축", "DV: 전환·리텐션·SUS·NASA-TLX·감정인지", "다음버전 우선순위 4"],
       "요건 ‘experiment design’ 강화 — 엄밀 통제실험.")

# SLIDE 19 — 비전 / SROI
s = slide(); bg(s, C["dusk"], C["ink"], angle=120)
random.seed(11)
for _ in range(30):
    rx, ry = random.uniform(0.3, 12.9), random.uniform(0.3, 3.0)
    d = random.choice([0.03, 0.05])
    o = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(rx), Inches(ry), Inches(d), Inches(d))
    o.fill.solid(); o.fill.fore_color.rgb = C["beige"]; o.line.fill.background(); o.shadow.inherit = False
txt(s, MARGIN, 1.0, CW, 0.4, [P([R("Ⅵ. 비전 — VISION", 13, C["amber"], True)], sa=0)])
txt(s, MARGIN, 1.5, CW, 1.4, [
    P([R("상담 이전과 일상 사이,", 30, C["white"], True)], sa=2),
    P([R("매일의 ‘마음 근육’을 키우는 예방적 동반자", 30, C["amber"], True)], sa=0)])
vis = [("예방적 케어", "병리가 아닌, 일상 속 자기관리 — 전문가 자문과 일치"),
       ("회기 간 공백을 채움", "상담과 상담 사이, 스스로 들여다보는 도구"),
       ("사회적 임팩트(SROI)", "상담 진입장벽 16.2% 현실에서 ‘이전 단계’를 메운다")]
for i, (t, d) in enumerate(vis):
    x = MARGIN + i*(CW/3)
    rect(s, x, 3.5, CW/3-0.3, 1.7, fill=rgb("#2E3149"), rounded=True, radius=0.08)
    txt(s, x+0.25, 3.7, CW/3-0.75, 1.4, [
        P([R(t, 15, C["mint"], True)], sa=6), P([R(d, 11.5, C["beige"])], sa=0, line=1.3)])
txt(s, MARGIN, 5.65, CW, 1.0, [
    P([R("“아주 보통의 하루에서도, 흘려보내던 나를 알아차리는 것.”", 18, C["white"], True)], sa=4),
    P([R("아보하는 거기서 시작합니다.", 14, C["beige"])], sa=0)], align=PP_ALIGN.CENTER)
note(s, "마지막으로 비전입니다. 아보하는 상담 이전과 일상 사이를 메우는 예방적 동반자를 지향합니다. 상담 전문가도 강조한 ‘회기 간 공백’을 "
        "채우고, 서비스 이용률 16%라는 현실에서 그 이전 단계를 담당합니다. 아주 보통의 하루에서도 흘려보내던 나를 알아차리는 것 — 아보하는 거기서 시작합니다.")
record(19, "비전 / SROI",
       ["예방적 케어(병리 아님)", "회기 간 공백을 채움", "SROI: 상담 이전 단계"],
       "비전·사회적 임팩트로 마무리(피드백 반영).")

# SLIDE 20 — Thank you
s = slide(); bg(s, C["coral"], rgb("#D24E3B"), angle=120)
pic_fit(s, asset("frontend", "public", "images", "mascot.png"), 10.0, 4.3, 2.6, 2.6, ha="center", va="bottom")
txt(s, MARGIN, 2.6, 10, 1.4, [
    P([R("감사합니다", 44, C["white"], True)], sa=4),
    P([R("아주 보통의 하루를, 가장 소중한 기록으로.", 18, C["dawn"])], sa=0)])
txt(s, MARGIN, 6.4, 10, 0.5, [P([R("Q & A   ·   닥토 공방 · 임동현", 14, C["white"], True)], sa=0)])
note(s, "이상으로 아보하 발표를 마칩니다. 질문 받겠습니다. 감사합니다.")
record(20, "감사합니다 / Q&A", ["Q&A"], "마무리.")

# ---------------------------------------------------------------------
prs.save(os.path.join(HERE, "아보하_HCI_최종발표.pptx"))
print("deck: done ->", os.path.join(HERE, "아보하_HCI_최종발표.pptx"))

# outline.md
with open(os.path.join(HERE, "outline.md"), "w", encoding="utf-8") as f:
    f.write("# 아보하(하루보석) — HCI 최종 발표 · 슬라이드 개요\n\n")
    f.write(f"총 {len(OUTLINE)}장 · 한국어 · 16:9 · 발표 ~12분 기준\n\n---\n\n")
    for n, title, bl, spk in OUTLINE:
        f.write(f"## {n}. {title}\n\n")
        for b in bl:
            f.write(f"- {b}\n")
        f.write(f"\n**발표 노트:** {spk}\n\n---\n\n")
print("outline: done ->", os.path.join(HERE, "outline.md"))
print("ALL DONE.")
